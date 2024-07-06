from pptx import Presentation


class PPTXParser:
    def get_pptx_text(pptx_files):
        text = ""
        for pptx_file in pptx_files:
            prs = Presentation(pptx_file)
            fullText = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        fullText.append(shape.text)
            text += '\n'.join(fullText)
        return text